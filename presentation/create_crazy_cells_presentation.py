#!/usr/bin/env python3

import csv
from pathlib import Path

from PIL import Image
from pptx import Presentation


ROOT = Path("/Users/finnborchers/Desktop/Cryo_Biokältetechnik")
TEMPLATE = ROOT / "presentation" / "Vorlage_Präsen.pptx"
OUTFILE = ROOT / "presentation" / "CrazyCells_Presentation.pptx"
FIG_DIR = ROOT / "analysis" / "figures"
PHOTO_DIR = ROOT / "presentation" / "images"
CONVERTED_PHOTO_DIR = PHOTO_DIR / "_converted"
NUCLEATION_EVENTS = FIG_DIR / "nucleation_events.csv"


def ensure_compatible_image(image_path: Path) -> Path:
    with Image.open(image_path) as im:
        fmt = (im.format or "").upper()
        if fmt in {"JPEG", "PNG", "GIF", "BMP", "TIFF", "WMF"}:
            return image_path
        CONVERTED_PHOTO_DIR.mkdir(parents=True, exist_ok=True)
        out = CONVERTED_PHOTO_DIR / f"{image_path.stem}.jpg"
        rgb = im.convert("RGB")
        rgb.save(out, format="JPEG", quality=95)
        return out


def remove_all_slides(prs: Presentation) -> None:
    # Remove template placeholder slides so we can rebuild while keeping the theme.
    for i in reversed(range(len(prs.slides))):
        slide_id = prs.slides._sldIdLst[i]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[i]


def set_notes(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


def set_bullets(placeholder, bullets):
    tf = placeholder.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0


def add_picture_fit(slide, image_path: Path, left, top, width, height):
    compatible_path = ensure_compatible_image(image_path)
    with Image.open(compatible_path) as im:
        iw, ih = im.size
    box_w = float(width)
    box_h = float(height)
    img_ratio = iw / ih
    box_ratio = box_w / box_h

    if img_ratio > box_ratio:
        draw_w = box_w
        draw_h = box_w / img_ratio
        draw_l = left
        draw_t = int(top + (box_h - draw_h) / 2)
    else:
        draw_h = box_h
        draw_w = box_h * img_ratio
        draw_t = top
        draw_l = int(left + (box_w - draw_w) / 2)

    slide.shapes.add_picture(
        str(compatible_path),
        draw_l,
        draw_t,
        width=int(draw_w),
        height=int(draw_h),
    )


def add_title_slide(prs, title, subtitle, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    set_notes(slide, notes)
    return slide


def add_title_content_slide(prs, title, bullets, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    set_bullets(body, bullets)
    set_notes(slide, notes)


def add_image_with_bullets_slide(prs, title, image_path, bullets, notes):
    # Layout 7: content with caption/body area.
    slide = prs.slides.add_slide(prs.slide_layouts[7])
    slide.shapes.title.text = title

    content = slide.placeholders[1]
    add_picture_fit(slide, image_path, content.left, content.top, content.width, content.height)

    body = slide.placeholders[2]
    set_bullets(body, bullets)

    set_notes(slide, notes)


def add_two_images_slide(prs, title, left_image, right_image, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two content
    slide.shapes.title.text = title

    left_ph = slide.placeholders[1]
    right_ph = slide.placeholders[2]

    add_picture_fit(slide, left_image, left_ph.left, left_ph.top, left_ph.width, left_ph.height)
    add_picture_fit(slide, right_image, right_ph.left, right_ph.top, right_ph.width, right_ph.height)

    set_notes(slide, notes)


def add_comparison_images_slide(
    prs,
    title,
    left_caption,
    left_image,
    right_caption,
    right_image,
    notes,
):
    # Layout 4: Comparison (caption + image on both sides).
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.shapes.title.text = title
    slide.placeholders[1].text = left_caption
    slide.placeholders[3].text = right_caption

    left_ph = slide.placeholders[2]
    right_ph = slide.placeholders[4]
    add_picture_fit(slide, left_image, left_ph.left, left_ph.top, left_ph.width, left_ph.height)
    add_picture_fit(slide, right_image, right_ph.left, right_ph.top, right_ph.width, right_ph.height)
    set_notes(slide, notes)


def add_cover_band_image(prs, slide, image_path: Path):
    # Add a wide image band to the lower part of title slide while keeping template title area.
    left = 0
    top = int(prs.slide_height * 0.42)
    width = int(prs.slide_width)
    height = int(prs.slide_height * 0.58)
    add_picture_fit(slide, image_path, left, top, width, height)


def nucleation_bullets(csv_path: Path):
    rows = []
    with csv_path.open("r", encoding="utf-8", newline="") as f:
        reader = csv.DictReader(f)
        for row in reader:
            rows.append(row)

    short = {
        "Crazy Cells-DMSO (°C)": "CC-DMSO",
        "Crazy Cells-Suc (°C)": "CC-Suc",
        "Crazy Cells-PBS (°C)": "CC-PBS",
        "Cryo Masters- DMSO (°C)": "CM-DMSO",
        "Cryo Masters- Suc (°C)": "CM-Suc",
        "Cryo Masters- PBS (°C)": "CM-PBS",
    }

    order = [
        "Crazy Cells-DMSO (°C)",
        "Crazy Cells-Suc (°C)",
        "Crazy Cells-PBS (°C)",
        "Cryo Masters- DMSO (°C)",
        "Cryo Masters- Suc (°C)",
        "Cryo Masters- PBS (°C)",
    ]
    rank = {name: i for i, name in enumerate(order)}
    ordered = sorted(rows, key=lambda r: rank.get(r["channel"], 999))
    bullets = ["Exact nucleation points (minimum -> rebound peak):"]
    for r in ordered:
        ch = short.get(r["channel"], r["channel"])
        bullets.append(
            f"{ch}: i{int(float(r['nucleation_index']))} ({float(r['nucleation_temp_c']):.2f} C) "
            f"-> i{int(float(r['peak_index']))} ({float(r['peak_temp_c']):.2f} C), "
            f"DeltaT {float(r['jump_delta_t_c']):.2f} C"
        )
    return bullets


def main():
    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)

    title_slide = add_title_slide(
        prs,
        "Cryopreservation of HeLa Cell Suspensions",
        "Crazy Cells (Finn Borchers, Franziska Reddner, Maria Shabanova)\nSupervisor: M.Sc. Tarek Deeb",
        (
            "Good afternoon everyone. We are the Crazy Cells group: Finn Borchers, Franziska Reddner, "
            "and Maria Shabanova. Our supervisor for this lab was M.Sc. Tarek Deeb. "
            "Today we present our cryopreservation project on HeLa cell suspensions. "
            "The central question is: how can we maximize survival after freezing and thawing by choosing "
            "the right cryoprotective condition and by controlling temperature over time?"
        ),
    )
    add_cover_band_image(prs, title_slide, PHOTO_DIR / "IMG_1788.jpeg")

    add_title_content_slide(
        prs,
        "Problem Statement: Why Cryopreservation Is Challenging",
        [
            "Cryopreservation is essential for cell therapy, research, and biobanking.",
            "Freezing can trigger lethal ice damage and osmotic (solution) damage.",
            "Goal: preserve viability and recovery after thawing.",
        ],
        (
            "Cryopreservation is widely used whenever living cells must be stored and transported. "
            "But freezing is not harmless. Cells are damaged mainly by two mechanisms: intracellular ice "
            "damage and osmotic or solution damage. So the practical problem is to keep cells viable and "
            "functional after thawing."
        ),
    )

    add_title_content_slide(
        prs,
        "Cryodamage Mechanisms and Working Hypothesis",
        [
            "2-factor concept: too slow cooling increases solution damage; too fast cooling increases ice damage.",
            "Cryoprotective agents (CPAs) reduce cryoinjury.",
            "Hypothesis: DMSO-based condition outperforms PBS-only after thawing.",
        ],
        (
            "Our background model follows the 2-factor hypothesis. If cooling is too slow, cells dehydrate too "
            "much and solution damage increases. If cooling is too fast, intracellular ice becomes likely. "
            "DMSO is a penetrating CPA and was expected to perform better than PBS-only."
        ),
    )

    add_comparison_images_slide(
        prs,
        "Experimental Setup in the Lab",
        "Sterile preparation and handling\n(Biosafety cabinet)",
        PHOTO_DIR / "IMG_1782.jpeg",
        "Cooling/measurement setup with probes\n(temperature logging during freezing)",
        PHOTO_DIR / "IMG_1784.jpeg",
        (
            "We split the HeLa suspension into three conditions: DMSO plus FBS, sucrose plus FBS, and PBS-only "
            "as low-protection reference. Samples were frozen, thawed, and measured in triplicates on Vi-CELL. "
            "Both teams generated data. We focus on Crazy Cells and use Cryo Masters for cross-checking trends."
        ),
    )

    add_image_with_bullets_slide(
        prs,
        "Data Sources and Analysis Metrics",
        PHOTO_DIR / "IMG_1789.jpeg",
        [
            "Temperature time series: 6 channels, 4,078 points each.",
            "Vi-CELL output: viability %, total cells, viable cells.",
            "Calculated: controlled cooling rate (5 to 45 min), cooling to nucleation, thawing, recovery.",
            "Recovery uses 1.47 million cells in 0.5 ml and 3 x 1 ml Vi-CELL aliquots after thaw.",
        ],
        (
            "We used two data sources. First, the temperature log for all team-condition combinations. "
            "Second, Vi-CELL measurements for viability and viable-cell concentration. "
            "From these we computed the required cooling, thawing, and recovery metrics. For recovery, each sample "
            "started from 0.5 ml before freezing and contained 1.47 million viable cells in total. After thawing, "
            "the cells were resuspended in 3 ml PBS and then split into three 1 ml aliquots for Vi-CELL counting."
        ),
    )

    add_image_with_bullets_slide(
        prs,
        "Temperature Profiles Across Freeze-Hold-Thaw",
        FIG_DIR / "01_temperature_profiles_full.png",
        [
            "All channels cool from ~6-7 C to near -190 C, then rewarm.",
            "Three phases are visible: cooling, deep-cold hold, thawing.",
            "Cryo Masters sucrose channel contained sensor outliers (cleaned).",
        ],
        (
            "This slide shows full temperature trajectories. We see clear cooling, deep-cold hold, and thawing "
            "phases. Overall behavior is consistent between teams and conditions. One channel, Cryo Masters "
            "sucrose, showed outlier spikes and was cleaned before final rate calculations."
        ),
    )

    add_two_images_slide(
        prs,
        "Cooling Performance vs Target Behavior",
        FIG_DIR / "02_temperature_profiles_zoom_freezing.png",
        FIG_DIR / "03_rate_summary.png",
        (
            "Zooming into the freezing window, we report controlled cooling from minute 5 to minute 45. "
            "These corrected rates are close to the target, around -0.93 to -0.96 K/min. "
            "Estimated nucleation temperatures ranged roughly from -12.7 to -7.1 C."
        ),
    )

    add_image_with_bullets_slide(
        prs,
        "Nucleation in Supercooled Water (Recalescence Zoom)",
        FIG_DIR / "06_nucleation_zoom.png",
        nucleation_bullets(NUCLEATION_EVENTS),
        (
            "This zoom shows the nucleation moment for each channel. Time zero is the supercooled minimum. "
            "After nucleation, temperature rebounds because latent heat is released as ice forms, and DeltaT "
            "captures the size of that jump. The largest jump here is Cryo Masters sucrose (DeltaT 10.90 C). "
            "These are thermal markers of freezing physics, not direct viability endpoints."
        ),
    )

    add_image_with_bullets_slide(
        prs,
        "Thawing Dynamics and Data Quality",
        FIG_DIR / "03_rate_summary.png",
        [
            "Thawing is much faster than controlled cooling.",
            "For PBS, Cryo Masters reaches +10 C earlier than Crazy Cells.",
            "Cryo Masters sucrose shows a sharp jump near +10 C (interpret with caution).",
        ],
        (
            "During thawing, temperatures rise much faster than during controlled cooling. "
            "For PBS, Cryo Masters reaches +10 C earlier than Crazy Cells in this dataset. "
            "Crazy Cells PBS: start -187.97 C at 52.35 min, +10 C at 55.479 min, rate 63.27 K/min. "
            "Cryo Masters PBS: start -189.51 C at 52.35 min, +10 C at 54.413 min, rate 96.73 K/min. "
            "Sucrose rates are 78.74 K/min (Crazy Cells) and 102.14 K/min (Cryo Masters), but the "
            "Cryo Masters sucrose trace has a sharp spike near +10 C, so this channel should be interpreted "
            "with caution."
        ),
    )

    add_image_with_bullets_slide(
        prs,
        "Post-thaw Viability by CPA Condition",
        FIG_DIR / "04_viability_by_solution.png",
        [
            "Crazy Cells medians: DMSO+FBS 85.39%, Sucrose+FBS 76.00%, PBS-only 60.00%.",
            "Cryo Masters shows similar DMSO/Sucrose tendency, but higher PBS variability.",
            "Triplicates here are technical Vi-CELL repeats from the same thawed sample.",
        ],
        (
            "For Crazy Cells, DMSO plus FBS gave the highest viability, sucrose plus FBS was intermediate, and "
            "PBS-only was lowest. Cryo Masters shows the same core pattern for DMSO and sucrose, while PBS-only "
            "has large spread. Because these are technical repeat counts from one thawed sample, the spread mainly "
            "reflects measurement scatter, mixing, or handling effects. Overall, this still supports the protective role of CPAs."
        ),
    )

    add_image_with_bullets_slide(
        prs,
        "Recovery Relative to Pre-freeze Viable Cells",
        FIG_DIR / "05_recovery_by_solution.png",
        [
            "Crazy Cells recovery medians: DMSO+FBS 18.72%, Sucrose+FBS 17.73%, PBS-only 1.48%.",
            "Cryo Masters recovery medians: DMSO+FBS 20.46%, Sucrose+FBS 14.29%, PBS-only 0.74%.",
            "Both teams: PBS-only is clearly lowest in recovery.",
        ],
        (
            "Recovery uses the lab-script definition of cells recovered divided by cells frozen. Each sample "
            "started from 0.5 ml before freezing and contained 1.47 million viable cells in total. After thawing, washing, and centrifugation, the cells were "
            "resuspended in 3 ml PBS. This suspension was divided into three 1 ml aliquots for Vi-CELL counting. "
            "These aliquots are technical repeat measurements of the same recovered sample. Recovery then compares "
            "the 1.47 million viable cells at the start with the cells estimated for the full recovered suspension after thaw. "
            "With that correction, DMSO remains highest, sucrose is intermediate, and PBS-only is clearly lowest."
        ),
    )

    add_title_content_slide(
        prs,
        "Discussion: What Did We Learn?",
        [
            "CPA choice strongly influences biological outcome.",
            "Temperature control was consistent, but survival still depended on formulation.",
            "DMSO benefit matches penetrating CPA mechanism.",
            "Sucrose gives partial protection; PBS-only is insufficient.",
        ],
        (
            "Our comparison suggests thermal control alone is not enough. Even with similar cooling profiles, "
            "biological outcome changed substantially by CPA condition. This is consistent with the mechanism: "
            "DMSO penetrates and protects against intracellular ice, while PBS-only lacks dedicated cryoprotection."
        ),
    )

    add_title_content_slide(
        prs,
        "Conclusions and Next Steps",
        [
            "In this HeLa run, DMSO+FBS performed best overall.",
            "Sucrose+FBS improved outcomes compared with PBS-only.",
            "Future work: larger n and tighter thaw-handling standardization.",
            "Thank you for your attention. Questions are welcome.",
        ],
        (
            "To conclude: DMSO plus FBS was the strongest condition in our data, and sucrose plus FBS still "
            "outperformed PBS-only. The temperature logs show reproducible freezing behavior, while viability and "
            "recovery emphasize the importance of CPA selection. Thank you, and special thanks to M.Sc. Tarek Deeb."
        ),
    )

    prs.save(str(OUTFILE))
    print(f"Created: {OUTFILE}")


if __name__ == "__main__":
    main()
